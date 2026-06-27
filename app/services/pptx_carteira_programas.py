import tempfile
import unicodedata
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = Path(__file__).resolve().parent.parent

LOGO_CAPA = BASE_DIR / "assets" / "logo2.png"
LOGO_SLIDE = BASE_DIR / "assets" / "logo.png"

ICON_PATH = Path("app/static/icons")
SLIDE_FINAL = BASE_DIR / "assets" / "slide_obrigado.png"

# ===== ESTILO INSTITUCIONAL =====
AZUL_TITULO = RGBColor(0, 91, 170)
AZUL_TABELA = RGBColor(0, 59, 143)
CINZA_SUBTITULO = RGBColor(100, 116, 139)
CINZA_RODAPE = RGBColor(100, 100, 100)
TEXTO_ESCURO = RGBColor(30, 41, 59)
LINHA_ZEBRA = RGBColor(245, 247, 250)
BRANCO = RGBColor(255, 255, 255)

PROGRAMAS_EXCLUIR = [
    "ASSESSORIA EM EDUCAÇÃO",
    "CERTIFICAÇÃO PROFISSIONAL",
    "DROPS EAD - PARCERIA EJA SESI",
]

def normalizar_programa_sql(alias="p"):
    return f"""
        CASE
            WHEN UPPER({alias}.nome_programa) IN (
                'RS QUALIFICAÇÃO',
                'RS QUALIFICAÇÃO RECOMEÇAR'
            )
            THEN 'RS QUALIFICAÇÃO'
            ELSE {alias}.nome_programa
        END
    """

FONTE_TITULO = "Arial Black"
FONTE_TEXTO = "Arial"

async def gerar_pptx_carteira_programas(request, filtros: dict):
    ano = filtros.get("ano") or 2026

    dados = await buscar_dados_carteira_programas(request, ano)
    detalhes = await buscar_detalhes_programas(request, ano)
    interlocutores = await buscar_programas_por_interlocutor(request, ano)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # limpa slides existentes do template
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    criar_slide_capa(prs, ano)
    criar_slide_tabela(prs, dados, ano)
    criar_slide_interlocutores(prs, interlocutores)

    for programa in detalhes:
        criar_slide_programa(prs, programa, ano)
    
    criar_slide_final(prs)

    saida = Path(tempfile.gettempdir()) / f"carteira_de_programas_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"

    for idx, slide in enumerate(prs.slides, start=1):
        if idx != 1 and idx != len(prs.slides):
            adicionar_numero_slide(slide, idx)

    prs.save(saida)

    return str(saida)


async def buscar_dados_carteira_programas(request, ano: int):
    pool = request.app.state.pool

    async with pool.acquire() as conn:
        rows = await conn.fetch("""
        SELECT
            CASE
                WHEN UPPER(p.nome_programa) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                THEN 'RS QUALIFICAÇÃO'
                ELSE p.nome_programa
            END AS nome_programa,

            string_agg(
                DISTINCT f.nome_financiamento,
                ' | '
                ORDER BY f.nome_financiamento
            ) AS tipo_financiamento,

            string_agg(
                DISTINCT m.nome,
                ', '
                ORDER BY m.nome
            ) AS modalidades,

            string_agg(
                DISTINCT COALESCE(i.nome, ''),
                ', '
                ORDER BY COALESCE(i.nome, '')
            ) AS interlocutor

        FROM ofertas_programas o

        LEFT JOIN programas p
            ON p.codigo = o.cod_programa

        LEFT JOIN financiamento f
            ON f.codigo = o.cod_financiamento

        LEFT JOIN modalidade m
            ON m.codigo = o.cod_modalidade

        LEFT JOIN interlocutores i
            ON i.codigo = p.cod_interlocutor

        WHERE o.ano = $1::int
            AND p.nome_programa IS NOT NULL
            AND TRIM(p.nome_programa) <> ''
            AND UPPER(p.nome_programa) NOT IN (
                'ASSESSORIA EM EDUCAÇÃO',
                'CERTIFICAÇÃO PROFISSIONAL',
                'DROPS EAD - PARCERIA EJA SESI'
            )

        GROUP BY
            CASE
                WHEN UPPER(p.nome_programa) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                THEN 'RS QUALIFICAÇÃO'
                ELSE p.nome_programa
            END

        ORDER BY
            nome_programa;
        """, int(ano))

    return [dict(r) for r in rows]

async def buscar_detalhes_programas(request, ano: int):
    pool = request.app.state.pool

    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            WITH base AS (
                SELECT
                    CASE
                        WHEN UPPER(p.nome_programa) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                        THEN 'RS QUALIFICAÇÃO'
                        ELSE p.nome_programa
                    END AS nome_programa,

                    MIN(p.codigo) AS cod_programa,
                    MAX(COALESCE(p.descricao, '')) AS objetivo,
                    string_agg(DISTINCT COALESCE(i.nome, ''), ', ' ORDER BY COALESCE(i.nome, '')) AS interlocutor,

                    COUNT(DISTINCT r.codigo) AS regioes_atendidas,
                    COUNT(DISTINCT o.cod_uo) AS unidades_atendidas,
                    COUNT(DISTINCT o.cod_modalidade) AS modalidades_qtd,

                    string_agg(DISTINCT r.nome, ', ' ORDER BY r.nome) AS regioes,
                    string_agg(DISTINCT m.nome, ', ' ORDER BY m.nome) AS modalidades,
                    string_agg(DISTINCT f.nome_financiamento, ', ' ORDER BY f.nome_financiamento) AS tipo_financiamento

                FROM ofertas_programas o

                LEFT JOIN programas p
                    ON p.codigo = o.cod_programa

                LEFT JOIN modalidade m
                    ON m.codigo = o.cod_modalidade
                
                LEFT JOIN financiamento f
                    ON f.codigo = o.cod_financiamento

                LEFT JOIN uo u
                    ON u.codigo::text = o.cod_uo::text

                LEFT JOIN subregioes s
                    ON s.codigo = u.cod_subregiao

                LEFT JOIN regioes r
                    ON r.codigo = s.codigo_regiao

                LEFT JOIN interlocutores i
                    ON i.codigo = p.cod_interlocutor

                WHERE o.ano = $1::int
                    AND p.nome_programa IS NOT NULL
                    AND TRIM(p.nome_programa) <> ''
                    AND UPPER(p.nome_programa) NOT IN (
                        'ASSESSORIA EM EDUCAÇÃO',
                        'CERTIFICAÇÃO PROFISSIONAL',
                        'DROPS EAD - PARCERIA EJA SESI'
                    )

                GROUP BY
                    CASE
                        WHEN UPPER(p.nome_programa) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                        THEN 'RS QUALIFICAÇÃO'
                        ELSE p.nome_programa
                    END
            ),

            realizado AS (
                SELECT
                    CASE
                        WHEN UPPER(p.nome_programa) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                        THEN 'RS QUALIFICAÇÃO'
                        ELSE p.nome_programa
                    END AS nome_programa,

                    SUM(COALESCE(vmr.matriculas_real, 0)) AS matriculas_realizado,
                    SUM(COALESCE(vmr.ha_real, 0)) AS ha_realizado,
                    SUM(COALESCE(vmr.receita_real, 0)) AS receita_realizado

                FROM ofertas_programas o

                LEFT JOIN programas p
                    ON p.codigo = o.cod_programa

                LEFT JOIN vw_meta_realizado vmr
                    ON vmr.cod_oferta = o.codigo
                   AND vmr.ano = o.ano

                WHERE o.ano = $1::int

                GROUP BY
                    CASE
                        WHEN UPPER(p.nome_programa) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                        THEN 'RS QUALIFICAÇÃO'
                        ELSE p.nome_programa
                    END
            ),

            metas AS (
                SELECT
                    CASE
                        WHEN UPPER(p.nome_programa) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                        THEN 'RS QUALIFICAÇÃO'
                        ELSE p.nome_programa
                    END AS nome_programa,

                    SUM(
                        CASE
                            WHEN mp.mes < EXTRACT(MONTH FROM CURRENT_DATE)
                            THEN COALESCE(mp.matriculas_meta, 0)
                            ELSE 0
                        END
                    ) AS matriculas_meta_atual,

                    SUM(COALESCE(mp.matriculas_meta, 0)) AS matriculas_meta_ano,

                    SUM(
                        CASE
                            WHEN mp.mes < EXTRACT(MONTH FROM CURRENT_DATE)
                            THEN COALESCE(mp.ha_meta, 0)
                            ELSE 0
                        END
                    ) AS ha_meta_atual,

                    SUM(COALESCE(mp.ha_meta, 0)) AS ha_meta_ano,

                    SUM(
                        CASE
                            WHEN mp.mes < EXTRACT(MONTH FROM CURRENT_DATE)
                            THEN COALESCE(mp.receita_meta, 0)
                            ELSE 0
                        END
                    ) AS receita_meta_atual,

                    SUM(COALESCE(mp.receita_meta, 0)) AS receita_meta_ano

                FROM ofertas_programas o

                LEFT JOIN programas p
                    ON p.codigo = o.cod_programa

                LEFT JOIN meta_programas mp
                    ON mp.cod_oferta = o.codigo
                   AND mp.ano = o.ano

                WHERE o.ano = $1::int

                GROUP BY
                    CASE
                        WHEN UPPER(p.nome_programa) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                        THEN 'RS QUALIFICAÇÃO'
                        ELSE p.nome_programa
                    END
            )

            SELECT
                b.cod_programa,
                b.nome_programa,
                b.objetivo,
                b.interlocutor,
                b.regioes_atendidas,
                b.unidades_atendidas,
                b.modalidades_qtd,
                b.regioes,
                b.modalidades,
                b.tipo_financiamento,

                COALESCE(r.matriculas_realizado, 0) AS matriculas_realizado,
                COALESCE(m.matriculas_meta_atual, 0) AS matriculas_meta_atual,
                COALESCE(m.matriculas_meta_ano, 0) AS matriculas_meta_ano,

                COALESCE(r.ha_realizado, 0) AS ha_realizado,
                COALESCE(m.ha_meta_atual, 0) AS ha_meta_atual,
                COALESCE(m.ha_meta_ano, 0) AS ha_meta_ano,

                COALESCE(r.receita_realizado, 0) AS receita_realizado,
                COALESCE(m.receita_meta_atual, 0) AS receita_meta_atual,
                COALESCE(m.receita_meta_ano, 0) AS receita_meta_ano

            FROM base b

            LEFT JOIN realizado r
                ON r.nome_programa = b.nome_programa

            LEFT JOIN metas m
                ON m.nome_programa = b.nome_programa

            ORDER BY
                b.nome_programa;
        """, int(ano))

    return [dict(r) for r in rows]

async def buscar_programas_por_interlocutor(request, ano: int):
    pool = request.app.state.pool

    async with pool.acquire() as conn:
        rows = await conn.fetch("""
            SELECT
                COALESCE(i.nome, 'Sem interlocutor') AS interlocutor,

                CASE
                    WHEN UPPER(TRIM(p.nome_programa)) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                    THEN 'RS QUALIFICAÇÃO'
                    ELSE TRIM(p.nome_programa)
                END AS nome_programa,

                CASE
                    WHEN f.nome_financiamento ILIKE '%GRATUIDADE NÃO REGIMENTAL%' THEN 'gnr'
                    WHEN f.nome_financiamento ILIKE '%GRATUIDADE REGIMENTAL%' THEN 'gr'
                    ELSE 'pago'
                END AS grupo

            FROM ofertas_programas o
            LEFT JOIN programas p ON p.codigo = o.cod_programa
            LEFT JOIN financiamento f ON f.codigo = o.cod_financiamento
            LEFT JOIN interlocutores i ON i.codigo = p.cod_interlocutor

            WHERE o.ano = $1::int
              AND p.nome_programa IS NOT NULL
              AND TRIM(p.nome_programa) <> ''
              AND UPPER(TRIM(p.nome_programa)) NOT IN (
                  'ASSESSORIA EM EDUCAÇÃO',
                  'CERTIFICAÇÃO PROFISSIONAL',
                  'DROPS EAD - PARCERIA EJA SESI'
              )

            GROUP BY
                COALESCE(i.nome, 'Sem interlocutor'),

                CASE
                    WHEN UPPER(TRIM(p.nome_programa)) IN ('RS QUALIFICAÇÃO', 'RS QUALIFICAÇÃO RECOMEÇAR')
                    THEN 'RS QUALIFICAÇÃO'
                    ELSE TRIM(p.nome_programa)
                END,

                CASE
                    WHEN f.nome_financiamento ILIKE '%GRATUIDADE NÃO REGIMENTAL%' THEN 'gnr'
                    WHEN f.nome_financiamento ILIKE '%GRATUIDADE REGIMENTAL%' THEN 'gr'
                    ELSE 'pago'
                END

            ORDER BY interlocutor, grupo, nome_programa;
        """, int(ano))

    return [dict(r) for r in rows]

def criar_slide_capa(prs, ano):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo branco/azulado
    fundo = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = RGBColor(248, 251, 255)
    fundo.line.fill.background()

    # Logo SENAI
    adicionar_logo(slide, capa=True)

    # Elemento azul grande inferior direito
    forma_azul = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.0),
        Inches(4.15),
        Inches(5.6),
        Inches(3.8)
    )
    forma_azul.fill.solid()
    forma_azul.fill.fore_color.rgb = AZUL_TABELA
    forma_azul.line.fill.background()
    forma_azul.rotation = -45

    # Linha diagonal decorativa
    linha = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(9.35),
        Inches(2.65),
        Inches(0.02),
        Inches(4.6)
    )
    linha.fill.solid()
    linha.fill.fore_color.rgb = RGBColor(96, 150, 255)
    linha.line.fill.background()
    linha.rotation = -45

    # Losangos claros
    losango1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.9),
        Inches(-0.75),
        Inches(2.3),
        Inches(2.3)
    )
    losango1.fill.solid()
    losango1.fill.fore_color.rgb = RGBColor(224, 234, 250)
    losango1.line.fill.background()
    losango1.rotation = 45

    losango2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0),
        Inches(5.95),
        Inches(2.1),
        Inches(2.1)
    )
    losango2.fill.solid()
    losango2.fill.fore_color.rgb = RGBColor(236, 242, 252)
    losango2.line.fill.background()
    losango2.rotation = 45

    # Pontilhado superior esquerdo
    for i in range(6):
        for j in range(6):
            ponto = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.32 + i * 0.18),
                Inches(0.35 + j * 0.18),
                Inches(0.025),
                Inches(0.025)
            )
            ponto.fill.solid()
            ponto.fill.fore_color.rgb = RGBColor(90, 140, 220)
            ponto.line.fill.background()

    # Título
    titulo = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(2.55),
        Inches(7.5),
        Inches(1.45)
    )
    tf = titulo.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "CARTEIRA DE"
    p.font.name = FONTE_TITULO
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = AZUL_TABELA

    p = tf.add_paragraph()
    p.text = "PROGRAMAS"
    p.font.name = FONTE_TITULO
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = AZUL_TABELA

    # Subtítulo
    subtitulo = slide.shapes.add_textbox(
        Inches(0.78),
        Inches(4.75),
        Inches(7.5),
        Inches(0.45)
    )
    p = subtitulo.text_frame.paragraphs[0]
    p.text = f"Portfólio institucional • Ano {ano}"
    p.font.name = FONTE_TEXTO
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(55, 65, 81)

    # Data com ícone
    add_icon(slide, "calendario.png", 0.78, 6.55, 0.45)

    data = slide.shapes.add_textbox(
        Inches(1.45),
        Inches(6.64),
        Inches(4.8),
        Inches(0.3)
    )
    p = data.text_frame.paragraphs[0]
    p.text = f"Gerado automaticamente em {datetime.now().strftime('%d/%m/%Y')}"
    p.font.name = FONTE_TEXTO
    p.font.size = Pt(10)
    p.font.color.rgb = CINZA_SUBTITULO


def classificar_financiamento(item):
    txt = (item.get("tipo_financiamento") or "").upper()

    if "PAGO POR PESSOA FÍSICA OU EMPRESA" in txt:
        return "pago"

    if "GRATUIDADE NÃO REGIMENTAL" in txt:
        return "gnr"

    if "GRATUIDADE REGIMENTAL" in txt:
        return "gr"

    return "pago"

def desenhar_tabela_programas_financiamento(slide, titulo, cor, itens, x, y, w, h):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(0.75)

    # Faixa de título
    faixa = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(0.55)
    )
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = cor
    faixa.line.fill.background()

    add_text(
        slide, titulo,
        x + 0.55, y + 0.15, w - 0.75, 0.25,
        11, BRANCO, True
    )

    linhas = len(itens)

    altura_tabela = 0.32 + (linhas * 0.28)

    tabela = slide.shapes.add_table(
        linhas + 1,
        3,
        Inches(x),
        Inches(y + 0.55),
        Inches(w),
        Inches(altura_tabela)
    ).table

    tabela.rows[0].height = Inches(0.32)

    larguras = [0.45, w - 2.25, 1.80]
    for i, largura in enumerate(larguras):
        tabela.columns[i].width = Inches(largura)

    headers = ["Nº", "Nome do Programa", "Interlocutor"]

    for c, texto in enumerate(headers):
        cell = tabela.cell(0, c)
        cell.text = texto.upper()
        cell.fill.solid()
        cell.fill.fore_color.rgb = cor
        cell.vertical_anchor = 3

        for p in cell.text_frame.paragraphs:
            p.font.name = FONTE_TEXTO
            p.font.size = Pt(6.8)
            p.font.bold = True
            p.font.color.rgb = BRANCO
            p.alignment = PP_ALIGN.CENTER

    for idx in range(1, linhas + 1):
        item = itens[idx - 1] if idx - 1 < len(itens) else {}

        row = idx
        tabela.rows[row].height = Inches(0.28)

        cor_linha = BRANCO if idx % 2 == 1 else RGBColor(245, 248, 252)

        valores = [
            f"{idx:02d}" if item else "",
            item.get("nome_programa", "") if item else "",
            item.get("interlocutor", "") if item else ""
        ]

        for c, valor in enumerate(valores):
            cell = tabela.cell(row, c)
            cell.text = str(valor or "")
            cell.fill.solid()
            cell.fill.fore_color.rgb = cor_linha
            cell.vertical_anchor = 3

            for p in cell.text_frame.paragraphs:
                p.font.name = FONTE_TEXTO
                p.font.size = Pt(6.7)
                p.font.color.rgb = TEXTO_ESCURO
                p.alignment = PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT

                if c == 1 and item:
                    p.font.bold = True

def desenhar_legenda_financiamento(slide, titulo, texto, cor, icone, x, y):
    add_icon(slide, icone, x, y - 0.04, 0.52)

    add_text(slide, titulo, x + 0.65, y + 0.02, 2.4, 0.20, 9, cor, True)
    add_text(slide, texto, x + 0.65, y + 0.28, 3.2, 0.35, 7.4, TEXTO_ESCURO)

def criar_slide_tabela(prs, dados, ano):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Cabeçalho
    barra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.75)
    )
    barra.fill.solid()
    barra.fill.fore_color.rgb = AZUL_TABELA
    barra.line.fill.background()

    adicionar_logo(slide)

    add_text(
        slide, "PROGRAMAS VIGENTES",
        0.35, 0.18, 8.5, 0.4,
        24, BRANCO, True, FONTE_TITULO
    )

    grupos = {
        "pago": {
            "titulo": "PAGO",
            "cor": AZUL_TABELA,
            "x": 0.25,
            "w": 4.25,
            "itens": []
        },
        "gr": {
            "titulo": "GRATUIDADE REGIMENTAL",
            "cor": RGBColor(24, 121, 45),
            "x": 4.65,
            "w": 4.25,
            "itens": []
        },
        "gnr": {
            "titulo": "GRATUIDADE NÃO REGIMENTAL",
            "cor": RGBColor(80, 35, 140),
            "x": 9.05,
            "w": 4.00,
            "itens": []
        }
    }

    for item in dados:
        grupos[classificar_financiamento(item)]["itens"].append(item)

    for grupo in grupos.values():
        desenhar_tabela_programas_financiamento(
            slide,
            titulo=grupo["titulo"],
            cor=grupo["cor"],
            itens=grupo["itens"],
            x=grupo["x"],
            y=1.05,
            w=grupo["w"],
            h=4.95
        )

    # Legenda inferior
    add_card(slide, 0.25, 6.18, 12.80, 0.75)

    desenhar_legenda_financiamento(
        slide, "PAGO",
        "Programas com cobrança de mensalidade ou contraprestação financeira.",
        AZUL_TABELA,
        "pg.png",
        0.65, 6.32
    )

    desenhar_legenda_financiamento(
        slide, "GRATUIDADE REGIMENTAL",
        "Programas ofertados gratuitamente conforme previsão da gratuidade regimental.",
        RGBColor(24, 121, 45),
        "gr.png",
        4.85, 6.32
    )

    desenhar_legenda_financiamento(
        slide, "GRATUIDADE NÃO REGIMENTAL",
        "Programas gratuitos não enquadrados na gratuidade regimental.",
        RGBColor(80, 35, 140),
        "gnr.png",
        9.15, 6.32
    )

def adicionar_logo(slide, capa=False):

    logo = LOGO_CAPA if capa else LOGO_SLIDE

    if not logo.exists():
        return

    if capa:
        slide.shapes.add_picture(
            str(logo),
            Inches(10.85),
            Inches(0.45),
            width=Inches(1.75)
        )
    else:
        slide.shapes.add_picture(
            str(logo),
            Inches(11.15),
            Inches(0.13),
            width=Inches(1.55)
        )

def criar_lista_texto(valor):
    if not valor:
        return "—"

    itens = [
        x.strip()
        for x in str(valor).split(",")
        if x.strip()
    ]

    return "\n".join([f"• {x}" for x in itens]) if itens else "—"

def add_text(slide, text, x, y, w, h, size=10, color=TEXTO_ESCURO, bold=False, font=FONTE_TEXTO):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = str(text or "")
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    return box


def add_card(slide, x, y, w, h):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BRANCO
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(0.75)
    return card


def add_icon(slide, arquivo, x, y, size=0.72):
    slide.shapes.add_picture(
        str(ICON_PATH / arquivo),
        Inches(x),
        Inches(y),
        width=Inches(size)
    )


def criar_card_indicador_exec(slide, titulo, icon, realizado, meta_atual, meta_ano, x, y, moeda=False):
    add_card(slide, x, y, 3.95, 1.35)
    add_icon(slide, icon, x + 0.22, y + 0.18)

    def fmt(v):
        v = float(v or 0)
        if moeda:
            return f"R$ {v:,.0f}".replace(",", "X").replace(".", ",").replace("X", ".")
        return f"{v:,.0f}".replace(",", ".")

    add_text(slide, "Realizado", x + 0.90, y + 0.74, 2.6, 0.25, 14, AZUL_TABELA, True, FONTE_TITULO)
    add_text(slide, fmt(realizado), x + 0.90, y + 0.50, 2.6, 0.30, 16, TEXTO_ESCURO, True)
    add_text(slide, "Realizado", x + 0.90, y + 0.78, 2.6, 0.30, 17, AZUL_TABELA, True, FONTE_TITULO)

    linha = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.90), Inches(y + 1.02), Inches(2.85), Inches(0.01)
    )
    linha.fill.solid()
    linha.fill.fore_color.rgb = RGBColor(226, 232, 240)
    linha.line.fill.background()

    add_text(slide, "Meta até o momento", x + 0.90, y + 1.10, 1.8, 0.18, 8, TEXTO_ESCURO)
    add_text(slide, fmt(meta_atual), x + 2.90, y + 1.10, 0.8, 0.18, 8, TEXTO_ESCURO)

    add_text(slide, "Meta anual", x + 0.90, y + 1.25, 1.8, 0.18, 8, TEXTO_ESCURO)
    add_text(slide, fmt(meta_ano), x + 2.90, y + 1.25, 0.8, 0.18, 8, TEXTO_ESCURO)

def adicionar_numero_slide(slide, numero):
    box = slide.shapes.add_textbox(
        Inches(12.55),
        Inches(7.08),
        Inches(0.45),
        Inches(0.25)
    )
    p = box.text_frame.paragraphs[0]
    p.text = str(numero)
    p.font.name = FONTE_TEXTO
    p.font.size = Pt(8)
    p.font.color.rgb = CINZA_SUBTITULO
    p.alignment = PP_ALIGN.RIGHT

def criar_slide_programa(prs, programa, ano):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Cabeçalho azul
    barra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, Inches(0.75)
    )
    barra.fill.solid()
    barra.fill.fore_color.rgb = AZUL_TABELA
    barra.line.fill.background()

    adicionar_logo(slide)

    add_text(
        slide,
        programa.get("nome_programa", "Programa"),
        0.45, 0.18, 10.4, 0.42,
        24, BRANCO, True, FONTE_TITULO
    )

    # OBJETIVO
    add_card(slide, 0.30, 0.86, 12.75, 1.03)
    add_icon(slide, "alvo.png", 0.47, 0.97, 0.58)

    add_text(slide, "OBJETIVO", 1.32, 1.02, 3, 0.25, 13, AZUL_TABELA, True)
    add_text(
        slide,
        programa.get("objetivo") or "Objetivo do programa não informado.",
        1.32, 1.25, 11.4, 0.72,
        10.5, TEXTO_ESCURO
    )

    # ABRANGÊNCIA
    add_card(slide, 0.30, 2.00, 12.75, 2.80)
    add_icon(slide, "globo.png", 0.47, 2.10, 0.58)

    add_text(slide, "ABRANGÊNCIA", 1.32, 2.18, 3, 0.25, 13, AZUL_TABELA, True)

    tx = slide.shapes.add_textbox(
        Inches(1.32),
        Inches(2.48),
        Inches(11.3),
        Inches(0.35)
    )

    tf = tx.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.font.name = FONTE_TEXTO
    p.font.size = Pt(10)
    p.font.color.rgb = TEXTO_ESCURO

    r = p.add_run()
    r.text = "O programa está presente em "

    r = p.add_run()
    r.text = f"{programa['regioes_atendidas']} regiões do Estado"
    r.font.bold = True

    r = p.add_run()
    r.text = ", abrangendo "

    r = p.add_run()
    r.text = f"{programa['unidades_atendidas']} unidades operacionais"
    r.font.bold = True

    r = p.add_run()
    r.text = " e ofertando "

    r = p.add_run()
    r.text = f"{programa['modalidades_qtd']} modalidade(s)"
    r.font.bold = True

    r = p.add_run()
    r.text = " de educação profissional."

    # Coluna regiões
    add_icon(slide, "mapa.png", 0.47, 2.82, 0.58)
    add_text(slide, "REGIÕES", 1.32, 2.98, 2.5, 0.25, 11, AZUL_TABELA, True)
    add_text(slide, criar_lista_texto(programa.get("regioes")), 1.32, 3.22, 2.70, 1.45, 8.1, TEXTO_ESCURO)

    # Divisor 1
    linha = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.45), Inches(2.95), Inches(0.01), Inches(1.65)
    )
    linha.fill.solid()
    linha.fill.fore_color.rgb = RGBColor(226, 232, 240)
    linha.line.fill.background()

    # Divisor 2
    linha = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.55), Inches(2.95), Inches(0.01), Inches(1.65)
    )
    linha.fill.solid()
    linha.fill.fore_color.rgb = RGBColor(226, 232, 240)
    linha.line.fill.background()

    # Coluna modalidades
    add_icon(slide, "livro.png", 4.70, 2.82, 0.58)

    add_text(
        slide,
        "MODALIDADES",
        5.55,
        2.98,
        2.6,
        0.25,
        11,
        AZUL_TABELA,
        True
    )

    add_text(
        slide,
        criar_lista_texto(programa.get("modalidades")),
        5.55,
        3.22,
        2.80,
        1.45,
        8.1,
        TEXTO_ESCURO
    )

    # Coluna tipo de financiamento
    add_icon(slide, "financiamento.png", 8.75, 2.82, 0.58)

    add_text(
        slide,
        "TIPO DE FINANCIAMENTO",
        9.60,
        2.98,
        3.0,
        0.25,
        11,
        AZUL_TABELA,
        True
    )

    add_text(
        slide,
        criar_lista_texto(programa.get("tipo_financiamento")),
        9.60,
        3.22,
        2.60,
        1.45,
        8.1,
        TEXTO_ESCURO
    )

    # INTERLOCUTOR
    add_card(slide, 0.30, 4.90, 12.75, 0.70)

    add_icon(slide, "usuario.png", 0.47, 4.92, 0.58)
    add_text(slide, "INTERLOCUTOR", 1.32, 5.08, 3.0, 0.22, 12, AZUL_TABELA, True)
    add_text(slide, programa.get("interlocutor") or "—", 1.32, 5.34, 6.0, 0.22, 9.5, TEXTO_ESCURO)

    # INDICADORES DO PROGRAMA
    add_text(
        slide,
        "INDICADORES DO PROGRAMA",
        0.30, 5.78, 4.0, 0.22,
        11, AZUL_TABELA, True
    )

    tabela = slide.shapes.add_table(
        4,
        4,
        Inches(0.30),
        Inches(6.08),
        Inches(12.75),
        Inches(1.02)
    ).table

    cabecalho = ["Indicador", "Realizado", "Meta até o momento", "Meta anual"]

    for c, texto in enumerate(cabecalho):
        cell = tabela.cell(0, c)
        cell.text = texto
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_TABELA
        cell.vertical_anchor = 3

        for p in cell.text_frame.paragraphs:
            p.font.name = FONTE_TEXTO
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = BRANCO
            p.alignment = PP_ALIGN.CENTER

    def fmt_num(v):
        return f"{float(v or 0):,.0f}".replace(",", ".")

    def fmt_moeda(v):
        return f"R$ {float(v or 0):,.0f}".replace(",", "X").replace(".", ",").replace("X", ".")

    linhas = [
        [
            "Matrículas",
            fmt_num(programa.get("matriculas_realizado")),
            fmt_num(programa.get("matriculas_meta_atual")),
            fmt_num(programa.get("matriculas_meta_ano")),
        ],
        [
            "Hora-aluno",
            fmt_num(programa.get("ha_realizado")),
            fmt_num(programa.get("ha_meta_atual")),
            fmt_num(programa.get("ha_meta_ano")),
        ],
        [
            "Receita",
            fmt_moeda(programa.get("receita_realizado")),
            fmt_moeda(programa.get("receita_meta_atual")),
            fmt_moeda(programa.get("receita_meta_ano")),
        ],
    ]

    for r, linha in enumerate(linhas, start=1):
        for c, texto in enumerate(linha):
            cell = tabela.cell(r, c)
            cell.text = texto
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRANCO

            for p in cell.text_frame.paragraphs:
                p.font.name = FONTE_TEXTO
                p.font.size = Pt(8)
                p.font.color.rgb = TEXTO_ESCURO
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

                if c == 0:
                    p.font.bold = True

def criar_card_numero(slide, titulo, valor, x, y):
    card = slide.shapes.add_shape(
        1,
        Inches(x),
        Inches(y),
        Inches(1.75),
        Inches(0.85)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LINHA_ZEBRA
    card.line.color.rgb = RGBColor(226, 232, 240)

    box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.12), Inches(1.5), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = titulo
    p.font.name = FONTE_TEXTO
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = CINZA_SUBTITULO

    box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.38), Inches(1.5), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = f"{int(valor or 0):,}".replace(",", ".")
    p.font.name = FONTE_TITULO
    p.font.size = Pt(18)
    p.font.color.rgb = AZUL_TITULO

def add_text_nowrap(slide, text, x, y, w, h, size=10, color=TEXTO_ESCURO, bold=False, font=FONTE_TEXTO):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = str(text or "")
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    return box

def criar_slide_interlocutores(prs, dados):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    barra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.75)
    )
    barra.fill.solid()
    barra.fill.fore_color.rgb = AZUL_TABELA
    barra.line.fill.background()

    adicionar_logo(slide)

    add_text(
        slide, "PROGRAMAS x INTERLOCUTORES",
        0.35, 0.18, 9.0, 0.42,
        25, BRANCO, True, FONTE_TITULO
    )

    def nome_programa_chave(nome):
        nome = str(nome or "").strip().upper()
        nome = " ".join(nome.split())

        nome_sem_acento = unicodedata.normalize("NFKD", nome)
        nome_sem_acento = "".join(c for c in nome_sem_acento if not unicodedata.combining(c))

        if nome_sem_acento in ("RS QUALIFICACAO", "RS QUALIFICACAO RECOMECAR"):
            return "RS QUALIFICACAO"

        if nome_sem_acento in ("SEJA PRO+", "SEJA PRÓ+"):
            return "SEJA PRO+"

        return nome_sem_acento

    agrupado = {}
    for item in dados:
        nome = item["interlocutor"]
        agrupado.setdefault(nome, {"pago": [], "gr": [], "gnr": []})
        programa_nome = item["nome_programa"]
        programa_chave = nome_programa_chave(programa_nome)

        ja_existe = any(
            nome_programa_chave(p) == programa_chave
            for p in agrupado[nome][item["grupo"]]
        )

        if not ja_existe:
            agrupado[nome][item["grupo"]].append(programa_nome)

    interlocutores = list(agrupado.items())[:5]

    cores = [
        RGBColor(24, 121, 45),
        AZUL_TABELA,
        RGBColor(0, 113, 145),
        RGBColor(80, 35, 140),
        RGBColor(204, 132, 0),
    ]

    largura_card = 2.45
    espacamento = 0.14
    x_inicial = 0.25
    y_card = 1.02

    for idx, (nome, grupos) in enumerate(interlocutores):
        x = x_inicial + idx * (largura_card + espacamento)
        cor = cores[idx % len(cores)]

        add_card(slide, x, y_card, largura_card, 5.12)

        topo = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y_card),
            Inches(largura_card), Inches(1.25)
        )
        topo.fill.solid()
        topo.fill.fore_color.rgb = cor
        topo.line.fill.background()

        programas_distintos = set()

        for grupo_lista in grupos.values():
            for programa_nome in grupo_lista:
                programas_distintos.add(nome_programa_chave(programa_nome))

        qtd = len(programas_distintos)

        add_text_nowrap(slide, nome.upper(), x + 0.16, y_card + 0.52, largura_card - 0.32, 0.22, 8.2, BRANCO, True)
        add_text(slide, f"{qtd} programa{'s' if qtd != 1 else ''}", x + 0.16, y_card + 0.82, largura_card - 0.32, 0.22, 7.5, BRANCO, True)

        y = y_card + 1.45

        def bloco(titulo, programas, icone, cor_bloco, y):
            if not programas:
                return y

            add_icon(slide, icone, x + 0.15, y, 0.34)
            add_text(slide, titulo, x + 0.52, y + 0.05, largura_card - 0.62, 0.22, 7.3, cor_bloco, True)

            y += 0.36

            for p_nome in programas:
                add_text(slide, f"• {p_nome}", x + 0.52, y, largura_card - 0.62, 0.14, 5.1, TEXTO_ESCURO)
                y += 0.135

            y += 0.07
            return y

        y = bloco("GRATUIDADE REGIMENTAL", grupos["gr"], "gr.png", RGBColor(24, 121, 45), y)
        y = bloco("GRATUIDADE NÃO REGIMENTAL", grupos["gnr"], "gnr.png", RGBColor(80, 35, 140), y)
        y = bloco("PAGO", grupos["pago"], "pg.png", AZUL_TABELA, y)

    add_card(slide, 0.25, 6.33, 12.80, 0.72)

    desenhar_legenda_financiamento(
        slide, "PAGO",
        "Programas com cobrança de mensalidade ou contraprestação financeira.",
        AZUL_TABELA,
        "pg.png",
        0.65, 6.43
    )

    desenhar_legenda_financiamento(
        slide, "GRATUIDADE REGIMENTAL",
        "Programas ofertados gratuitamente conforme previsão da gratuidade regimental.",
        RGBColor(24, 121, 45),
        "gr.png",
        4.85, 6.43
    )

    desenhar_legenda_financiamento(
        slide, "GRATUIDADE NÃO REGIMENTAL",
        "Programas gratuitos não enquadrados na gratuidade regimental.",
        RGBColor(80, 35, 140),
        "gnr.png",
        9.15, 6.43
    )

def criar_slide_final(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if SLIDE_FINAL.exists():
        slide.shapes.add_picture(
            str(SLIDE_FINAL),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height
        )